"""Document text extraction utilities.

Extracts readable text from binary document formats so LLMs can
consume their content.  Each ``extract_*`` function takes a file path
and returns the document's text as a plain string.

All libraries used are permissively licensed (MIT / BSD).
"""

import zipfile


def extract_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx(file_path: str) -> str:
    """Extract text from a Word (.docx) file."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    parts = []
    for para in doc.paragraphs:
        parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def extract_xlsx(file_path: str) -> str:
    """Extract text from an Excel (.xlsx) file."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"--- {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            parts.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return "\n".join(parts)


def extract_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint (.pptx) file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def extract_rtf(file_path: str) -> str:
    """Extract text from a Rich Text Format (.rtf) file."""
    from striprtf.striprtf import rtf_to_text

    with open(file_path, "rb") as f:
        raw = f.read()
    return rtf_to_text(raw.decode("utf-8", errors="replace"))


def extract_xls(file_path: str) -> str:
    """Extract text from a legacy Excel (.xls) file."""
    import xlrd

    wb = xlrd.open_workbook(file_path)
    parts = []
    for sheet in wb.sheets():
        parts.append(f"--- {sheet.name} ---")
        for row_idx in range(sheet.nrows):
            parts.append("\t".join(
                str(sheet.cell_value(row_idx, col_idx))
                for col_idx in range(sheet.ncols)
            ))
    return "\n".join(parts)


def extract_odt(file_path: str) -> str:
    """Extract text from an OpenDocument Text (.odt) file."""
    from lxml import etree

    with zipfile.ZipFile(file_path) as zf:
        with zf.open("content.xml") as f:
            tree = etree.parse(f)
    ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    return "\n".join(
        "".join(p.itertext())
        for p in tree.iter(f"{{{ns}}}p")
    )


def extract_ods(file_path: str) -> str:
    """Extract text from an OpenDocument Spreadsheet (.ods) file."""
    from lxml import etree

    with zipfile.ZipFile(file_path) as zf:
        with zf.open("content.xml") as f:
            tree = etree.parse(f)
    ns_table = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
    ns_text = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    parts = []
    for table in tree.iter(f"{{{ns_table}}}table"):
        name = table.get(f"{{{ns_table}}}name", "Sheet")
        parts.append(f"--- {name} ---")
        for row in table.iter(f"{{{ns_table}}}table-row"):
            cells = []
            for cell in row.iter(f"{{{ns_table}}}table-cell"):
                cell_text = " ".join(
                    "".join(p.itertext())
                    for p in cell.iter(f"{{{ns_text}}}p")
                )
                cells.append(cell_text)
            parts.append("\t".join(cells))
    return "\n".join(parts)


def extract_odp(file_path: str) -> str:
    """Extract text from an OpenDocument Presentation (.odp) file."""
    from lxml import etree

    with zipfile.ZipFile(file_path) as zf:
        with zf.open("content.xml") as f:
            tree = etree.parse(f)
    ns_draw = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    ns_text = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    parts = []
    for i, page in enumerate(tree.iter(f"{{{ns_draw}}}page"), 1):
        parts.append(f"--- Slide {i} ---")
        for p in page.iter(f"{{{ns_text}}}p"):
            text = "".join(p.itertext()).strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def extract_epub(file_path: str) -> str:
    """Extract text from an EPUB e-book."""
    from lxml import etree

    parts = []
    with zipfile.ZipFile(file_path) as zf:
        # Parse the container to find the root file
        with zf.open("META-INF/container.xml") as cf:
            container = etree.parse(cf)
        ns_container = "urn:oasis:names:tc:opendocument:xmlns:container"
        rootfile = container.find(f".//{{{ns_container}}}rootfile")
        if rootfile is None:
            rootfile = container.xpath("//*[local-name()='rootfile']")
            rootfile = rootfile[0] if rootfile else None

        if rootfile is not None:
            opf_path = rootfile.get("full-path", "")
            opf_dir = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""
            with zf.open(opf_path) as opf_file:
                opf = etree.parse(opf_file)
            spine_ids = [
                item.get("idref")
                for item in opf.xpath("//*[local-name()='itemref']")
            ]
            manifest = {
                item.get("id"): item.get("href")
                for item in opf.xpath("//*[local-name()='item']")
            }
            for idref in spine_ids:
                href = manifest.get(idref, "")
                item_path = opf_dir + href if not href.startswith("/") else href.lstrip("/")
                try:
                    with zf.open(item_path) as html_file:
                        html_tree = etree.parse(html_file, etree.HTMLParser())
                        body = html_tree.find(".//body")
                        if body is not None:
                            text = "".join(body.itertext())
                            parts.append(text.strip())
                except (KeyError, etree.XMLSyntaxError):
                    continue
        else:
            for name in zf.namelist():
                if name.endswith((".html", ".xhtml", ".htm")):
                    try:
                        with zf.open(name) as html_file:
                            html_tree = etree.parse(html_file, etree.HTMLParser())
                            body = html_tree.find(".//body")
                            if body is not None:
                                text = "".join(body.itertext())
                                parts.append(text.strip())
                    except etree.XMLSyntaxError:
                        continue
    return "\n\n".join(parts)


def extract_eml(file_path: str) -> str:
    """Extract text from an email message (.eml)."""
    import email
    from email import policy

    with open(file_path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=policy.default)
    parts = []
    for header in ("From", "To", "Cc", "Date", "Subject"):
        val = msg.get(header)
        if val:
            parts.append(f"{header}: {val}")
    parts.append("")  # blank line after headers
    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            from lxml import etree
            tree = etree.HTML(content)
            content = "".join(tree.itertext()) if tree is not None else content
        parts.append(content)
    return "\n".join(parts)


# MIME type / extension → extractor mapping.
# Checked in order by read_file; the first match wins.
# Each entry: (mime_type_or_None, file_extension_or_None, extractor)
EXTRACTORS: list[tuple[str | None, str | None, callable]] = [
    ("application/pdf", None, extract_pdf),
    ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", None, extract_docx),
    ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", None, extract_xlsx),
    ("application/vnd.openxmlformats-officedocument.presentationml.presentation", None, extract_pptx),
    ("application/rtf", ".rtf", extract_rtf),
    ("application/vnd.ms-excel", ".xls", extract_xls),
    ("application/vnd.oasis.opendocument.text", ".odt", extract_odt),
    ("application/vnd.oasis.opendocument.spreadsheet", ".ods", extract_ods),
    ("application/vnd.oasis.opendocument.presentation", ".odp", extract_odp),
    ("application/epub+zip", ".epub", extract_epub),
    ("message/rfc822", ".eml", extract_eml),
]
